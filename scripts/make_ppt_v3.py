"""Regulatory Tracker — rich 3-slide pitch deck."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]

# ── palette ──────────────────────────────────────────────────────
DARK       = RGBColor(0x10, 0x14, 0x18)
DARK2      = RGBColor(0x1E, 0x24, 0x2B)
IBM_BLUE   = RGBColor(0x00, 0x43, 0xCE)
BLUE_LIGHT = RGBColor(0x00, 0x8A, 0xD8)
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
OFFWHITE   = RGBColor(0xF4, 0xF4, 0xF4)
MID        = RGBColor(0x69, 0x72, 0x7B)
RED        = RGBColor(0xDA, 0x1E, 0x28)
RED_LIGHT  = RGBColor(0xFF, 0xF0, 0xF0)
GREEN      = RGBColor(0x19, 0x8A, 0x38)
GREEN_PALE = RGBColor(0xDE, 0xF7, 0xE9)
PURPLE     = RGBColor(0x6F, 0x29, 0xC1)
AMBER      = RGBColor(0xF1, 0xC2, 0x1B)
AMBER_PALE = RGBColor(0xFF, 0xF8, 0xE1)

W = Inches(13.33)
H = Inches(7.5)


# ── helpers ──────────────────────────────────────────────────────
def rect(slide, x, y, w, h, fill, alpha=None):
    sh = slide.shapes.add_shape(1, x, y, w, h)
    sh.fill.solid()
    sh.fill.fore_color.rgb = fill
    sh.line.fill.background()
    return sh


def txt(slide, t, x, y, w, h, sz=16, bold=False, color=DARK,
        align=PP_ALIGN.LEFT, italic=False, wrap=True):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = t
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def multiline(slide, lines, x, y, w, h, sz=15, color=DARK,
              leading=0.6, bold_first=False):
    """Each element of lines is either a str or (str, bold, color)."""
    for i, line in enumerate(lines):
        if isinstance(line, tuple):
            t, bld, clr = line
        else:
            t, bld, clr = line, (bold_first and i == 0), color
        tb = slide.shapes.add_textbox(x, y + Inches(i * leading), w, Inches(leading + 0.05))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = t
        r.font.size = Pt(sz)
        r.font.bold = bld
        r.font.color.rgb = clr


def card(slide, x, y, w, h, header_color, header_text, body_lines,
         body_bg=OFFWHITE, header_sz=14, body_sz=13):
    rect(slide, x, y, w, Inches(0.48), header_color)
    txt(slide, header_text, x + Inches(0.12), y + Inches(0.04),
        w - Inches(0.24), Inches(0.42), sz=header_sz,
        bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(slide, x, y + Inches(0.48), w, h - Inches(0.48), body_bg)
    multiline(slide, body_lines,
              x + Inches(0.15), y + Inches(0.58),
              w - Inches(0.3), h - Inches(0.58),
              sz=body_sz, leading=0.52)


def stat_block(slide, x, y, w, h, number, label, bg, num_color=WHITE):
    rect(slide, x, y, w, h, bg)
    txt(slide, number, x, y + Inches(0.15), w, Inches(0.7),
        sz=34, bold=True, color=num_color, align=PP_ALIGN.CENTER)
    txt(slide, label, x, y + Inches(0.85), w, Inches(0.45),
        sz=12, color=num_color, align=PP_ALIGN.CENTER)


def tag(slide, x, y, label, bg, fg=WHITE, sz=11):
    rect(slide, x, y, Inches(1.5), Inches(0.32), bg)
    txt(slide, label, x, y, Inches(1.5), Inches(0.32),
        sz=sz, bold=True, color=fg, align=PP_ALIGN.CENTER)


def divider(slide, x, y, w, color=IBM_BLUE, h=0.04):
    rect(slide, x, y, w, Inches(h), color)


def slide_num(slide, n, total=3):
    txt(slide, f"{n:02d} / {total:02d}",
        Inches(12.6), Inches(7.1), Inches(0.65), Inches(0.32),
        sz=10, color=MID, align=PP_ALIGN.RIGHT)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 1 — THE PROBLEM
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, WHITE)

# ── left dark panel ──
rect(s, 0, 0, Inches(4.6), H, DARK)
divider(s, 0, 0, Inches(4.6), IBM_BLUE, h=0.07)

txt(s, "THE", Inches(0.4), Inches(0.5), Inches(3.8), Inches(0.55),
    sz=13, color=MID, bold=True)
txt(s, "Problem", Inches(0.4), Inches(0.95), Inches(4.0), Inches(1.5),
    sz=52, bold=True, color=WHITE)
divider(s, Inches(0.4), Inches(2.5), Inches(1.2), BLUE_LIGHT)
txt(s, "Hong Kong\nFinancial\nCompliance", Inches(0.4), Inches(2.7),
    Inches(3.8), Inches(1.6), sz=17, color=MID, italic=True)

# stat blocks at bottom of left panel
stat_block(s, Inches(0.25), Inches(4.8), Inches(1.95), Inches(1.55),
           "3", "Regulators", DARK2, WHITE)
stat_block(s, Inches(2.4), Inches(4.8), Inches(1.95), Inches(1.55),
           "100s", "of PDFs", IBM_BLUE, WHITE)

# ── right content ──
rect(s, Inches(4.6), 0, Inches(0.07), H, IBM_BLUE)

txt(s, "Every compliance question costs hours of manual work",
    Inches(4.9), Inches(0.3), Inches(8.1), Inches(0.65),
    sz=20, bold=True, color=DARK)
divider(s, Inches(4.9), Inches(1.02), Inches(7.8), MID, h=0.025)

# pain point cards — 2x2 grid
pain = [
    (IBM_BLUE,  "No Unified Search",
     ["3 separate regulator websites", "No cross-source query possible", "No alerts when rules change"]),
    (RED,       "Silent Rule Changes",
     ["Documents updated without notice", "No version tracking or diff", "Easy to miss critical updates"]),
    (PURPLE,    "Manual & Error-Prone",
     ["Hours spent reading PDFs", "Knowledge locked in individuals", "Can't scale with team growth"]),
    (RGBColor(0xB4,0x5A,0x00), "Real Regulatory Risk",
     ["Missed circular = potential breach", "Fines, licence suspension", "Reputational damage"]),
]
cw, ch = Inches(3.8), Inches(2.45)
positions = [
    (Inches(4.9), Inches(1.15)),
    (Inches(8.9), Inches(1.15)),
    (Inches(4.9), Inches(3.75)),
    (Inches(8.9), Inches(3.75)),
]
for (x, y), (color, header, body) in zip(positions, pain):
    card(s, x, y, cw, ch, color, header,
         [f"·  {b}" for b in body], body_bg=OFFWHITE, body_sz=13)

# bottom callout
rect(s, Inches(4.9), Inches(6.35), Inches(8.0), Inches(0.82), RED_LIGHT)
rect(s, Inches(4.9), Inches(6.35), Inches(0.08), Inches(0.82), RED)
txt(s, "\"I don't know if that circular was updated last week.\"  — Every compliance officer.",
    Inches(5.1), Inches(6.45), Inches(7.7), Inches(0.6),
    sz=14, italic=True, color=RED)

slide_num(s, 1)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 2 — THE PRODUCT
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, DARK)
divider(s, 0, 0, W, IBM_BLUE, h=0.07)
divider(s, 0, H - Inches(0.07), W, IBM_BLUE, h=0.07)

# ── hero ──
txt(s, "INTRODUCING", Inches(0.5), Inches(0.22), Inches(12.0), Inches(0.45),
    sz=12, color=MID, bold=True, align=PP_ALIGN.CENTER)
txt(s, "Regulatory Tracker", Inches(0.5), Inches(0.65), Inches(12.3), Inches(1.3),
    sz=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
txt(s, "Ask a compliance question in plain English. Get a precise, cited answer in seconds.",
    Inches(0.5), Inches(1.9), Inches(12.3), Inches(0.6),
    sz=18, color=MID, italic=True, align=PP_ALIGN.CENTER)
divider(s, Inches(5.5), Inches(2.6), Inches(2.3), BLUE_LIGHT)

# ── 3 feature columns ──
features = [
    (IBM_BLUE,  "ONE PLACE",
     "SFC · Insurance Authority · PCPD",
     ["Search across all 3 regulators simultaneously",
      "Covers Codes, Guidelines, Circulars & Guidance",
      "Always in sync as new documents are published"]),
    (PURPLE,    "PLAIN ENGLISH",
     "No boolean queries. No manual filtering.",
     ["Ask naturally: \"What are the AML requirements?\"",
      "Filter by regulator or search across all",
      "Understands regulatory context and terminology"]),
    (GREEN,     "ALWAYS CITED",
     "Every answer backed by the source.",
     ["Exact document, section, and regulator cited",
      "No hallucination — grounded in real documents",
      "Audit-ready output for compliance records"]),
]
cw = Inches(3.9)
for i, (color, title, subtitle, bullets_) in enumerate(features):
    x = Inches(0.45) + i * (cw + Inches(0.23))
    # top colour band
    rect(s, x, Inches(2.85), cw, Inches(0.06), color)
    rect(s, x, Inches(2.91), cw, Inches(3.95), DARK2)
    txt(s, title, x + Inches(0.2), Inches(3.0), cw - Inches(0.4), Inches(0.5),
        sz=16, bold=True, color=color)
    txt(s, subtitle, x + Inches(0.2), Inches(3.52), cw - Inches(0.4), Inches(0.45),
        sz=12, italic=True, color=MID)
    divider(s, x + Inches(0.2), Inches(4.0), cw - Inches(0.4), MID, h=0.02)
    multiline(s, [f"→  {b}" for b in bullets_],
              x + Inches(0.2), Inches(4.1), cw - Inches(0.4), Inches(2.5),
              sz=13, color=RGBColor(0xC6, 0xC6, 0xC6), leading=0.6)

# ── how it works strip ──
rect(s, 0, Inches(7.05), W, Inches(0.37), RGBColor(0x1E,0x24,0x2B))
steps = ["Collect Docs", "Index Content", "You Ask", "AI Answers", "Cite Source"]
colors = [IBM_BLUE, BLUE_LIGHT, PURPLE, GREEN, AMBER]
sw = Inches(13.33 / len(steps))
for i, (step, color) in enumerate(zip(steps, colors)):
    rect(s, i * sw, Inches(7.05), sw - Inches(0.04), Inches(0.37), color)
    txt(s, step, i * sw, Inches(7.08), sw, Inches(0.3),
        sz=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)

slide_num(s, 2)


# ═══════════════════════════════════════════════════════════════════
# SLIDE 3 — WHERE WE ARE
# ═══════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, WHITE)

# ── top dark header ──
rect(s, 0, 0, W, Inches(1.4), DARK)
divider(s, 0, 0, W, IBM_BLUE, h=0.07)
txt(s, "WHERE WE ARE", Inches(0.5), Inches(0.12), Inches(6.0), Inches(0.4),
    sz=11, bold=True, color=MID)
txt(s, "Progress & Roadmap", Inches(0.5), Inches(0.48), Inches(7.0), Inches(0.85),
    sz=36, bold=True, color=WHITE)
# tag — active branch
rect(s, Inches(10.0), Inches(0.5), Inches(2.9), Inches(0.5), IBM_BLUE)
txt(s, "Branch: fix/regulator-client-bugs",
    Inches(10.05), Inches(0.52), Inches(2.8), Inches(0.45),
    sz=11, color=WHITE, align=PP_ALIGN.CENTER)

# ── stat row ──
stats = [
    ("3",    "Regulators\nConnected",  IBM_BLUE),
    ("73+",  "Documents\nTracked",     PURPLE),
    ("100%", "Test Pass\nRate",        GREEN),
    ("0",    "Manual Steps\nRequired", DARK),
]
sw2 = Inches(3.1)
for i, (num, lbl, color) in enumerate(stats):
    x = Inches(0.35) + i * (sw2 + Inches(0.12))
    stat_block(s, x, Inches(1.55), sw2, Inches(1.4), num, lbl, color, WHITE)

# ── done column ──
rect(s, Inches(0.35), Inches(3.1), Inches(6.1), Inches(3.95), GREEN_PALE)
rect(s, Inches(0.35), Inches(3.1), Inches(0.08), Inches(3.95), GREEN)
txt(s, "✓  COMPLETED", Inches(0.55), Inches(3.18), Inches(5.7), Inches(0.45),
    sz=13, bold=True, color=GREEN)

done_items = [
    ("SFC scraper",        "16 Codes + 51 Guidelines discovered"),
    ("Insurance Authority","18 docs (2026) · 42 docs (2025) · full year sweep"),
    ("PCPD scraper",       "4 key privacy documents — idempotent"),
    ("Q&A engine",         "Hybrid search + reranking + answer generation"),
    ("Test suite",         "7 / 7 tests passing · CI green"),
]
for i, (label, detail) in enumerate(done_items):
    y = Inches(3.75) + i * Inches(0.62)
    rect(s, Inches(0.55), y, Inches(5.7), Inches(0.55), WHITE)
    txt(s, label, Inches(0.68), y + Inches(0.04), Inches(2.1), Inches(0.47),
        sz=13, bold=True, color=DARK)
    txt(s, detail, Inches(2.85), y + Inches(0.04), Inches(3.3), Inches(0.47),
        sz=12, color=MID)

# ── next column ──
rect(s, Inches(6.7), Inches(3.1), Inches(6.3), Inches(3.95), OFFWHITE)
rect(s, Inches(6.7), Inches(3.1), Inches(0.08), Inches(3.95), IBM_BLUE)
txt(s, "→  UP NEXT", Inches(6.9), Inches(3.18), Inches(5.9), Inches(0.45),
    sz=13, bold=True, color=IBM_BLUE)

next_items = [
    ("Read PDFs",       "Extract & index all document content",         "1"),
    ("Beta Launch",     "Open Q&A to compliance team for testing",      "2"),
    ("Go Live",         "Production deployment — real-time updates",    "3"),
]
for i, (label, detail, num) in enumerate(next_items):
    y = Inches(3.75) + i * Inches(0.9)
    rect(s, Inches(6.9), y, Inches(0.38), Inches(0.38), IBM_BLUE)
    txt(s, num, Inches(6.9), y, Inches(0.38), Inches(0.38),
        sz=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    txt(s, label, Inches(7.4), y, Inches(2.2), Inches(0.38),
        sz=14, bold=True, color=DARK)
    txt(s, detail, Inches(7.4), y + Inches(0.4), Inches(5.3), Inches(0.38),
        sz=12, color=MID)

# vision line
rect(s, 0, Inches(7.08), W, Inches(0.42), DARK)
txt(s, "Vision:  No compliance officer should ever manually read 50 PDFs to answer one question again.",
    Inches(0.5), Inches(7.1), Inches(12.3), Inches(0.35),
    sz=13, italic=True, color=MID, align=PP_ALIGN.CENTER)

slide_num(s, 3)


# ── save ──────────────────────────────────────────────────────────
out = "regulatory_tracker_v3.pptx"
prs.save(out)
print(f"Saved: {out}")
