"""Generate Regulatory Tracker Agent progress presentation as a .pptx file."""
from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ── palette ──────────────────────────────────────────────────────────────────
IBM_BLUE   = RGBColor(0x00, 0x43, 0xCE)   # IBM Blue 60
IBM_DARK   = RGBColor(0x16, 0x1A, 0x1E)   # near-black
WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF4, 0xF4, 0xF4)
MID_GRAY   = RGBColor(0x69, 0x72, 0x7B)
ACCENT     = RGBColor(0x00, 0x8A, 0xD8)   # lighter IBM blue for accents
GREEN      = RGBColor(0x19, 0x8A, 0x38)
YELLOW     = RGBColor(0xF1, 0xC2, 0x1B)

SLIDE_W = Inches(13.33)
SLIDE_H = Inches(7.5)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H

BLANK = prs.slide_layouts[6]   # completely blank


def add_rect(slide, left, top, width, height, color, alpha=None):
    shape = slide.shapes.add_shape(1, left, top, width, height)
    shape.line.fill.background()
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    return shape


def add_text(slide, text, left, top, width, height,
             font_size=18, bold=False, color=IBM_DARK,
             align=PP_ALIGN.LEFT, wrap=True, italic=False):
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    return txBox


def add_bullet_box(slide, items, left, top, width, height,
                   font_size=16, title=None, title_color=IBM_BLUE,
                   bullet="▸  ", bg_color=None, text_color=IBM_DARK):
    if bg_color:
        add_rect(slide, left, top, width, height, bg_color)
    if title:
        add_text(slide, title, left + Inches(0.15), top + Inches(0.1),
                 width - Inches(0.3), Inches(0.45),
                 font_size=font_size + 2, bold=True, color=title_color)
        top += Inches(0.5)
        height -= Inches(0.5)
    txBox = slide.shapes.add_textbox(left + Inches(0.15), top + Inches(0.1),
                                     width - Inches(0.3), height - Inches(0.1))
    tf = txBox.text_frame
    tf.word_wrap = True
    for i, item in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = f"{bullet}{item}"
        run.font.size = Pt(font_size)
        run.font.color.rgb = text_color


def header_bar(slide, title, subtitle=None):
    """Dark top bar with title."""
    add_rect(slide, 0, 0, SLIDE_W, Inches(1.25), IBM_DARK)
    add_text(slide, title,
             Inches(0.5), Inches(0.15), Inches(12), Inches(0.65),
             font_size=28, bold=True, color=WHITE)
    if subtitle:
        add_text(slide, subtitle,
                 Inches(0.5), Inches(0.82), Inches(12), Inches(0.4),
                 font_size=16, color=ACCENT)
    # thin blue accent line under header
    add_rect(slide, 0, Inches(1.25), SLIDE_W, Inches(0.06), IBM_BLUE)


# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 1 — Title
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, IBM_DARK)
add_rect(s, 0, Inches(5.5), SLIDE_W, Inches(0.08), IBM_BLUE)

add_text(s, "Regulatory Tracker Agent",
         Inches(1), Inches(1.6), Inches(11.3), Inches(1.2),
         font_size=48, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, "Hong Kong Financial Regulations · RAG System",
         Inches(1), Inches(2.85), Inches(11.3), Inches(0.6),
         font_size=24, color=ACCENT, align=PP_ALIGN.CENTER)
add_text(s, "SFC  ·  Insurance Authority  ·  PCPD",
         Inches(1), Inches(3.55), Inches(11.3), Inches(0.5),
         font_size=18, color=MID_GRAY, align=PP_ALIGN.CENTER)
add_text(s, "Progress Briefing  —  2025",
         Inches(1), Inches(5.7), Inches(11.3), Inches(0.45),
         font_size=14, color=MID_GRAY, align=PP_ALIGN.CENTER, italic=True)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 2 — Problem Statement
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
header_bar(s, "Problem Statement")

problems = [
    "3 Hong Kong financial regulators publish PDFs across separate, static HTML pages — no unified search API exists",
    "SFC, IA, and PCPD each use a different page structure: attribute-gated rows, year-paginated listings, and curated fixed sets",
    "Documents are republished with cache-busting query strings — naive URL hashing causes silent duplicate ingestion",
    "No way to answer cross-regulator compliance questions without manually reading dozens of PDFs",
]
add_bullet_box(s, problems,
               Inches(0.5), Inches(1.5), Inches(12.3), Inches(5.5),
               font_size=17, bg_color=LIGHT_GRAY, bullet="●  ")

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 3 — Solution Overview
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
header_bar(s, "Solution Overview", "End-to-end ingestion + retrieval pipeline")

steps = [
    ("1  DISCOVER", "Scrape SFC / IA / PCPD listing pages; emit normalised doc records", IBM_BLUE),
    ("2  STORE",    "Persist metadata + PDF bytes to IBM Cloud Object Storage (COS)",     ACCENT),
    ("3  CHUNK",    "Markdown-aware section splitter + sliding-window fallback",           RGBColor(0x6F,0x29,0xC1)),
    ("4  EMBED",    "WatsonX ibm/slate-30m-english-rtrvr-v2 → 1536-dim vectors",         IBM_BLUE),
    ("5  RETRIEVE", "Hybrid ANN (vector) + keyword search in AstraDB, merged via RRF",    ACCENT),
    ("6  RERANK",   "WatsonX cross-encoder reranker narrows candidates to top-K",         RGBColor(0x6F,0x29,0xC1)),
    ("7  GENERATE", "IBM Granite LLM produces a grounded, cited answer via WatsonX",      GREEN),
]

col_w = Inches(1.72)
gap   = Inches(0.06)
y0    = Inches(1.5)
for i, (label, desc, color) in enumerate(steps):
    x = Inches(0.3) + i * (col_w + gap)
    add_rect(s, x, y0, col_w, Inches(0.52), color)
    add_text(s, label, x, y0, col_w, Inches(0.52),
             font_size=11, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, x, y0 + Inches(0.52), col_w, Inches(4.6), LIGHT_GRAY)
    add_text(s, desc, x + Inches(0.08), y0 + Inches(0.6),
             col_w - Inches(0.16), Inches(4.4),
             font_size=12, color=IBM_DARK)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 4 — Regulator Coverage
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
header_bar(s, "Regulator Coverage", "3 PRs shipped and passing")

sources = [
    ("PR 1 — PCPD", GREEN,
     ["Curated fixed set — 4 privacy documents",
      "PDPO Full Ordinance",
      "Six Data Protection Principles — Overview",
      "Data Security Measures Guidance",
      "Data Breach Handling Guidance Note",
      "Tests: 2 / 2 ✓  |  Idempotency verified"]),
    ("PR 2 — Insurance Authority", IBM_BLUE,
     ["Year-paginated circular listing",
      "URL pattern: /circulars_on_regulatory_matters_{year}.html",
      "2026: 18 documents  |  2025: 42 documents",
      "Zero-results guard distinguishes 'no history' vs parser failure",
      "Full sweep from earliest year to current",
      "Tests: 2 / 2 ✓"]),
    ("PR 3 — SFC", ACCENT,
     ["Two pages: Codes (16 docs) + Guidelines (51 docs)",
      "Parser filters by data-code-guideline-id attribute only",
      "Handbook popup: 5 PDFs extracted from #popuphb block",
      "Previous-version popup rows correctly excluded",
      "HTML fixture regression test included",
      "Tests: 3 / 3 ✓"]),
]

col_w = Inches(3.9)
for i, (title, color, bullets) in enumerate(sources):
    x = Inches(0.35) + i * (col_w + Inches(0.25))
    add_rect(s, x, Inches(1.45), col_w, Inches(0.55), color)
    add_text(s, title, x, Inches(1.45), col_w, Inches(0.55),
             font_size=15, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, x, Inches(2.0), col_w, Inches(5.0), LIGHT_GRAY)
    add_bullet_box(s, bullets, x, Inches(2.0), col_w, Inches(5.0),
                   font_size=14, bullet="▸  ", text_color=IBM_DARK)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 5 — Ingestion Architecture
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
header_bar(s, "Ingestion Pipeline Architecture")

boxes = [
    (Inches(0.4),  Inches(2.2), "SFC Client\ncore/sfc_client.py",     ACCENT),
    (Inches(0.4),  Inches(3.6), "IA Client\ncore/ia_client.py",       IBM_BLUE),
    (Inches(0.4),  Inches(5.0), "PCPD Client\ncore/pcpd_client.py",   RGBColor(0x6F,0x29,0xC1)),
    (Inches(4.0),  Inches(3.3), "doc_id hash\nsha1(url_path)[:12]",   IBM_DARK),
    (Inches(7.5),  Inches(2.4), "IBM COS\nPDF bytes\npdfs/<id>.pdf",  ACCENT),
    (Inches(7.5),  Inches(4.3), "AstraDB\nDoc metadata\nJSON records", IBM_BLUE),
    (Inches(11.0), Inches(3.3), "POST /ingest/bulk\nFastAPI",          GREEN),
]
for (x, y, label, color) in boxes:
    add_rect(s, x, y, Inches(2.0), Inches(1.0), color)
    add_text(s, label, x + Inches(0.08), y + Inches(0.05),
             Inches(1.84), Inches(0.9),
             font_size=12, bold=False, color=WHITE)

# arrows (thin blue lines as rectangles)
arrow_data = [
    # clients → doc_id box
    (Inches(2.4),  Inches(2.7),  Inches(1.6),  Inches(0.05)),
    (Inches(2.4),  Inches(4.1),  Inches(1.6),  Inches(0.05)),
    (Inches(2.4),  Inches(5.5),  Inches(1.6),  Inches(0.05)),
    # doc_id → COS
    (Inches(6.0),  Inches(2.9),  Inches(1.5),  Inches(0.05)),
    # doc_id → AstraDB
    (Inches(6.0),  Inches(4.8),  Inches(1.5),  Inches(0.05)),
    # COS + Astra → API
    (Inches(9.5),  Inches(2.9),  Inches(1.5),  Inches(0.05)),
    (Inches(9.5),  Inches(4.8),  Inches(1.5),  Inches(0.05)),
]
for (x, y, w, h) in arrow_data:
    add_rect(s, x, y, w, h, MID_GRAY)

add_text(s, "discover_documents() → normalised record list",
         Inches(0.4), Inches(6.4), Inches(8.0), Inches(0.4),
         font_size=13, color=MID_GRAY, italic=True)
add_text(s, "doc_id = f\"{source}-{sha1(url_path)[:12]}\"",
         Inches(4.0), Inches(6.4), Inches(5.0), Inches(0.4),
         font_size=13, color=IBM_BLUE, bold=True)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 6 — RAG Pipeline Architecture
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
header_bar(s, "RAG Pipeline Architecture", "Hybrid retrieval → rerank → generate")

rag_steps = [
    ("Chunker\ncore/chunker.py",         RGBColor(0x6F,0x29,0xC1), "Markdown heading split\n+ sliding window"),
    ("Embedder\ncore/embedder.py",        IBM_BLUE,                 "WatsonX Slate 30M\n1536-dim vectors"),
    ("AstraDB\nchunks collection",        ACCENT,                   "ANN vector search\n+ keyword (BM25)"),
    ("RRF Merge\ncore/retriever.py",      IBM_DARK,                 "Reciprocal Rank\nFusion top-N"),
    ("Reranker\ncore/reranker.py",        IBM_BLUE,                 "WatsonX cross-encoder\nor local fallback"),
    ("Generator\ncore/generator.py",      GREEN,                    "IBM Granite LLM\ncitation-backed answer"),
]

bw = Inches(1.95)
gap = Inches(0.1)
y_top = Inches(1.55)
for i, (label, color, note) in enumerate(rag_steps):
    x = Inches(0.35) + i * (bw + gap)
    add_rect(s, x, y_top, bw, Inches(0.85), color)
    add_text(s, label, x + Inches(0.06), y_top + Inches(0.05),
             bw - Inches(0.12), Inches(0.78),
             font_size=12, bold=True, color=WHITE)
    add_rect(s, x, y_top + Inches(0.85), bw, Inches(1.2), LIGHT_GRAY)
    add_text(s, note, x + Inches(0.06), y_top + Inches(0.9),
             bw - Inches(0.12), Inches(1.1),
             font_size=12, color=IBM_DARK)
    if i < len(rag_steps) - 1:
        ax = x + bw + Inches(0.01)
        add_text(s, "→", ax, y_top + Inches(0.25), gap + Inches(0.05), Inches(0.45),
                 font_size=18, bold=True, color=IBM_BLUE, align=PP_ALIGN.CENTER)

# Chat endpoint box
add_rect(s, Inches(0.35), Inches(4.0), Inches(12.6), Inches(1.0), LIGHT_GRAY)
add_text(s, "POST /chat  →  ChatRequest { question, source_filter, top_k }  →  ChatResponse { answer, citations[], model_used, chunk_count }",
         Inches(0.5), Inches(4.05), Inches(12.3), Inches(0.9),
         font_size=14, color=IBM_DARK, bold=False)

add_text(s, "Source filter: restrict retrieval to SFC | IA | PCPD, or omit for cross-regulator search",
         Inches(0.5), Inches(5.2), Inches(12.0), Inches(0.45),
         font_size=13, color=MID_GRAY, italic=True)
add_text(s, "Citations embedded inline — format: [Source: <source>, <section>, doc: <doc_id>]",
         Inches(0.5), Inches(5.7), Inches(12.0), Inches(0.45),
         font_size=13, color=MID_GRAY, italic=True)
add_text(s, "IAM token cached 50 min  ·  Batch embed (10 texts/call)  ·  WatsonX or local reranker (RERANKER env var)",
         Inches(0.5), Inches(6.2), Inches(12.0), Inches(0.45),
         font_size=13, color=MID_GRAY, italic=True)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 7 — API Surface
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
header_bar(s, "API Surface", "FastAPI · uvicorn · port 8000")

# left col — ingest
add_rect(s, Inches(0.4), Inches(1.45), Inches(5.9), Inches(5.6), LIGHT_GRAY)
add_text(s, "POST /ingest/bulk",
         Inches(0.55), Inches(1.5), Inches(5.6), Inches(0.5),
         font_size=18, bold=True, color=IBM_BLUE)
ingest_items = [
    "Query param: source = SFC | IA | PCPD",
    "Calls discover_documents() on the selected client",
    "Deduplicates via doc_id (SHA-1 of URL path)",
    "Persists new records to COS or in-memory store",
    "Returns: { ingested: n, skipped: n, source: str }",
    "Idempotent — safe to run repeatedly",
]
add_bullet_box(s, ingest_items, Inches(0.4), Inches(2.05),
               Inches(5.9), Inches(4.8), font_size=14)

# right col — chat
add_rect(s, Inches(6.8), Inches(1.45), Inches(6.1), Inches(5.6), LIGHT_GRAY)
add_text(s, "POST /chat",
         Inches(6.95), Inches(1.5), Inches(5.8), Inches(0.5),
         font_size=18, bold=True, color=IBM_BLUE)
chat_items = [
    "Body: { question, source_filter?, top_k (1–20) }",
    "Hybrid retrieval: ANN vector + keyword search",
    "RRF merge → WatsonX rerank → top-K chunks",
    "IBM Granite generates grounded answer",
    "Returns: { answer, citations[], model_used, chunk_count }",
    "source_filter: restrict to SFC / IA / PCPD or omit for all",
]
add_bullet_box(s, chat_items, Inches(6.8), Inches(2.05),
               Inches(6.1), Inches(4.8), font_size=14)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 8 — Tech Stack
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
header_bar(s, "Tech Stack")

categories = [
    ("API & Framework",  IBM_BLUE,  ["FastAPI + Uvicorn", "Pydantic v2 schemas", "Python 3.10+"]),
    ("IBM AI / Watson",  ACCENT,    ["WatsonX ibm/granite-13b-chat-v2 (LLM)", "WatsonX ibm/slate-30m-english-rtrvr-v2 (embeddings)", "WatsonX cross-encoder reranker"]),
    ("Storage",          RGBColor(0x6F,0x29,0xC1), ["AstraDB (DataStax) — vector + keyword chunks", "IBM Cloud Object Storage — PDFs + metadata JSON", "astrapy client + ibm-cos-sdk"]),
    ("Scraping",         GREEN,     ["BeautifulSoup4 — HTML parsing", "requests — HTTP client", "0.5–1s politeness delay between requests"]),
    ("Testing & CI",     IBM_DARK,  ["pytest + httpx", "GitHub Actions CI (lint, tests, commit-lint)", "HTML fixture regression tests for SFC parser"]),
]

cw = Inches(2.35)
gap = Inches(0.12)
for i, (cat, color, items) in enumerate(categories):
    x = Inches(0.35) + i * (cw + gap)
    add_rect(s, x, Inches(1.45), cw, Inches(0.52), color)
    add_text(s, cat, x, Inches(1.45), cw, Inches(0.52),
             font_size=13, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, x, Inches(1.97), cw, Inches(5.1), LIGHT_GRAY)
    add_bullet_box(s, items, x, Inches(1.97), cw, Inches(5.1),
                   font_size=12, bullet="·  ", text_color=IBM_DARK)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 9 — Current Branch & Status
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
header_bar(s, "Current Status", "Branch: fix/regulator-client-bugs")

# status table header
add_rect(s, Inches(0.4), Inches(1.5), Inches(12.5), Inches(0.5), IBM_DARK)
for x, w, label in [
    (Inches(0.4),  Inches(1.5), "PR"),
    (Inches(1.9),  Inches(2.0), "Source"),
    (Inches(3.9),  Inches(5.5), "What it does"),
    (Inches(9.4),  Inches(1.6), "Tests"),
    (Inches(11.0), Inches(1.9), "Live count"),
]:
    add_text(s, label, x + Inches(0.08), Inches(1.52), w, Inches(0.46),
             font_size=14, bold=True, color=WHITE)

rows = [
    ("PR 1", "PCPD",  "Curated list of 4 privacy documents",                   "2 / 2 ✓", "4 docs"),
    ("PR 2", "IA",    "Scrapes IA circulars by year; year-range sweep",         "2 / 2 ✓", "18 (2026)  42 (2025)"),
    ("PR 3", "SFC",   "Codes + Guidelines; Handbook popup; prev-version guard", "3 / 3 ✓", "16 Codes  51 Guidelines"),
]
for i, (pr, src, desc, tests, count) in enumerate(rows):
    bg = LIGHT_GRAY if i % 2 == 0 else WHITE
    y = Inches(2.0) + i * Inches(0.75)
    add_rect(s, Inches(0.4), y, Inches(12.5), Inches(0.72), bg)
    for x, w, val, clr, bld in [
        (Inches(0.4),  Inches(1.5), pr,    IBM_DARK, True),
        (Inches(1.9),  Inches(2.0), src,   IBM_BLUE, True),
        (Inches(3.9),  Inches(5.5), desc,  IBM_DARK, False),
        (Inches(9.4),  Inches(1.6), tests, GREEN,    True),
        (Inches(11.0), Inches(1.9), count, IBM_DARK, False),
    ]:
        add_text(s, val, x + Inches(0.08), y + Inches(0.1), w, Inches(0.55),
                 font_size=13, bold=bld, color=clr)

# Active branch note
add_rect(s, Inches(0.4), Inches(4.6), Inches(12.5), Inches(0.7), RGBColor(0xFF,0xF8,0xE1))
add_text(s, "⚙  Active branch:  fix/regulator-client-bugs  — bug fixes on regulator scraper clients currently in progress",
         Inches(0.55), Inches(4.65), Inches(12.2), Inches(0.6),
         font_size=14, color=IBM_DARK)

# new files
add_text(s, "New files added (untracked in git):",
         Inches(0.5), Inches(5.5), Inches(12.0), Inches(0.35),
         font_size=13, bold=True, color=IBM_DARK)
new_files = "core/chunker.py  ·  core/embedder.py  ·  core/generator.py  ·  core/reranker.py  ·  core/retriever.py  ·  api/routers/chat.py  ·  store/astra_chunk_store.py  ·  store/astra_layout_store.py  ·  store/cos_document_store.py"
add_text(s, new_files,
         Inches(0.5), Inches(5.9), Inches(12.3), Inches(0.7),
         font_size=12, color=MID_GRAY)

# ═══════════════════════════════════════════════════════════════════════════
# SLIDE 10 — Next Steps
# ═══════════════════════════════════════════════════════════════════════════
s = prs.slides.add_slide(BLANK)
add_rect(s, 0, 0, SLIDE_W, SLIDE_H, WHITE)
header_bar(s, "Next Steps")

next_steps = [
    ("PDF Download & OCR",       IBM_BLUE,  ["Download raw PDF bytes for all ingested doc records", "Run Docling to convert PDFs → Markdown", "Store results in COS (transformed/<id>.md)"]),
    ("Layout Store",             ACCENT,    ["Populate AstraDB astra_layout_store with bounding-box elements", "Enable page-level citation lookup for the chat API", "Wire store/astra_layout_store.py into pipeline"]),
    ("End-to-End Chunk Pipeline", RGBColor(0x6F,0x29,0xC1), ["Run scripts/run_chunk.py on all downloaded documents", "Embed all chunks via WatsonX Slate model", "Populate AstraDB chunks collection for RAG retrieval"]),
    ("Production Hardening",     IBM_DARK,  ["Merge fix/regulator-client-bugs branch", "Add rate-limit retry + backoff to all HTTP clients", "Production env var secrets via IBM Key Protect"]),
]

cw = Inches(2.95)
gap = Inches(0.15)
for i, (title, color, items) in enumerate(next_steps):
    x = Inches(0.35) + i * (cw + gap)
    add_rect(s, x, Inches(1.45), cw, Inches(0.52), color)
    add_text(s, title, x, Inches(1.45), cw, Inches(0.52),
             font_size=14, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    add_rect(s, x, Inches(1.97), cw, Inches(4.9), LIGHT_GRAY)
    add_bullet_box(s, items, x, Inches(1.97), cw, Inches(4.9),
                   font_size=13, bullet="▸  ", text_color=IBM_DARK)

# closing note
add_text(s, "Goal: fully populated AstraDB chunks collection → /chat endpoint live against real regulatory content",
         Inches(0.5), Inches(7.0), Inches(12.3), Inches(0.4),
         font_size=14, bold=True, color=IBM_BLUE, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════
# SAVE
# ═══════════════════════════════════════════════════════════════════════════
out = "regulatory_tracker_progress.pptx"
prs.save(out)
print(f"Saved: {out}")
