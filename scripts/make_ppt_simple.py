"""Generate simple 3-slide Regulatory Tracker pitch deck."""
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN

prs = Presentation()
prs.slide_width  = Inches(13.33)
prs.slide_height = Inches(7.5)

BLANK = prs.slide_layouts[6]

IBM_DARK  = RGBColor(0x16, 0x1A, 0x1E)
IBM_BLUE  = RGBColor(0x00, 0x43, 0xCE)
WHITE     = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT     = RGBColor(0xF4, 0xF4, 0xF4)
MID       = RGBColor(0x69, 0x72, 0x7B)
RED       = RGBColor(0xDA, 0x1E, 0x28)
GREEN     = RGBColor(0x19, 0x8A, 0x38)

W = Inches(13.33)
H = Inches(7.5)


def rect(slide, x, y, w, h, color):
    s = slide.shapes.add_shape(1, x, y, w, h)
    s.fill.solid()
    s.fill.fore_color.rgb = color
    s.line.fill.background()
    return s


def text(slide, txt, x, y, w, h, size=18, bold=False, color=IBM_DARK,
         align=PP_ALIGN.LEFT, italic=False):
    tb = slide.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = txt
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color
    return tb


def bullets(slide, items, x, y, w, h, size=17, color=IBM_DARK, spacing=0.72):
    for i, item in enumerate(items):
        tb = slide.shapes.add_textbox(x, y + Inches(i * spacing), w, Inches(spacing))
        tf = tb.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        r = p.add_run()
        r.text = f"—   {item}"
        r.font.size = Pt(size)
        r.font.color.rgb = color


# ─────────────────────────────────────────────────────────────────
# SLIDE 1 — The Problem
# ─────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, WHITE)

# left dark panel
rect(s, 0, 0, Inches(5.2), H, IBM_DARK)
text(s, "The\nProblem", Inches(0.45), Inches(2.2), Inches(4.5), Inches(2.8),
     size=52, bold=True, color=WHITE)
text(s, "Hong Kong Financial Compliance", Inches(0.45), Inches(5.3),
     Inches(4.5), Inches(0.5), size=13, color=MID, italic=True)

# accent bar
rect(s, Inches(5.2), 0, Inches(0.08), H, IBM_BLUE)

# right content
text(s, "Every compliance question starts the same way:",
     Inches(5.6), Inches(1.1), Inches(7.3), Inches(0.6),
     size=17, color=MID, italic=True)

pain_points = [
    "Open 3 regulator websites  (SFC, IA, PCPD)",
    "Search manually through hundreds of PDFs",
    "Hope nothing changed since you last checked",
    "Repeat for every new question",
]
bullets(s, pain_points, Inches(5.6), Inches(1.9), Inches(7.2), Inches(3.0),
        size=18, color=IBM_DARK, spacing=0.78)

# callout box
rect(s, Inches(5.6), Inches(5.15), Inches(7.2), Inches(1.65), RGBColor(0xFF,0xF0,0xF0))
rect(s, Inches(5.6), Inches(5.15), Inches(0.1), Inches(1.65), RED)
text(s, "One missed circular = a regulatory breach.\nFines. Licence suspensions. Reputational damage.",
     Inches(5.85), Inches(5.25), Inches(6.8), Inches(1.4),
     size=16, bold=True, color=RED)

# slide number
text(s, "01 / 03", Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.35),
     size=11, color=MID, align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────────────────────────
# SLIDE 2 — The Product
# ─────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, IBM_DARK)

# top label
rect(s, 0, 0, W, Inches(0.08), IBM_BLUE)
text(s, "THE PRODUCT", Inches(0.5), Inches(0.2), Inches(12.0), Inches(0.5),
     size=13, color=MID, align=PP_ALIGN.CENTER)

# hero line
text(s, "Regulatory Tracker", Inches(0.5), Inches(0.85), Inches(12.3), Inches(1.4),
     size=54, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
text(s, "Ask a compliance question. Get the answer — with the exact source cited.",
     Inches(0.5), Inches(2.25), Inches(12.3), Inches(0.65),
     size=20, color=MID, align=PP_ALIGN.CENTER, italic=True)

# three feature cards
cards = [
    (IBM_BLUE,                  "One Place",       "All 3 Hong Kong regulators — SFC, IA, and PCPD — in a single search."),
    (RGBColor(0x6F,0x29,0xC1), "Plain English",    "Ask in natural language. No boolean queries, no manual filtering."),
    (GREEN,                     "Always Cited",    "Every answer references the exact document, section, and regulator."),
]
cw = Inches(3.8)
for i, (color, title, desc) in enumerate(cards):
    x = Inches(0.55) + i * (cw + Inches(0.28))
    rect(s, x, Inches(3.2), cw, Inches(0.7), color)
    text(s, title, x, Inches(3.2), cw, Inches(0.7),
         size=18, bold=True, color=WHITE, align=PP_ALIGN.CENTER)
    rect(s, x, Inches(3.9), cw, Inches(2.5), RGBColor(0x26, 0x2B, 0x30))
    text(s, desc, x + Inches(0.15), Inches(4.05),
         cw - Inches(0.3), Inches(2.2),
         size=15, color=RGBColor(0xC6, 0xC6, 0xC6))

text(s, "02 / 03", Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.35),
     size=11, color=MID, align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────────────────────────
# SLIDE 3 — Where We Are
# ─────────────────────────────────────────────────────────────────
s = prs.slides.add_slide(BLANK)
rect(s, 0, 0, W, H, WHITE)

# left dark panel
rect(s, 0, 0, Inches(5.2), H, IBM_DARK)
text(s, "Where\nWe Are", Inches(0.45), Inches(2.0), Inches(4.5), Inches(2.8),
     size=52, bold=True, color=WHITE)
rect(s, Inches(5.2), 0, Inches(0.08), H, IBM_BLUE)

# done section
text(s, "✓  Done", Inches(5.6), Inches(0.85), Inches(7.2), Inches(0.55),
     size=17, bold=True, color=GREEN)
done = [
    "All 3 regulators connected and tested",
    "73+ regulatory documents discovered and tracked",
    "Q&A engine built — retrieval, reranking, and answer generation",
]
bullets(s, done, Inches(5.6), Inches(1.5), Inches(7.2), Inches(2.4),
        size=16, color=IBM_DARK, spacing=0.68)

# divider
rect(s, Inches(5.6), Inches(3.7), Inches(7.0), Inches(0.04), LIGHT)

# next section
text(s, "→  Next", Inches(5.6), Inches(3.9), Inches(7.2), Inches(0.55),
     size=17, bold=True, color=IBM_BLUE)
next_steps = [
    "Read & index all PDF content",
    "Open beta with compliance team",
    "Production launch",
]
bullets(s, next_steps, Inches(5.6), Inches(4.55), Inches(7.2), Inches(2.2),
        size=16, color=IBM_DARK, spacing=0.68)

text(s, "03 / 03", Inches(12.5), Inches(7.1), Inches(0.7), Inches(0.35),
     size=11, color=MID, align=PP_ALIGN.RIGHT)


# ─────────────────────────────────────────────────────────────────
# SAVE
# ─────────────────────────────────────────────────────────────────
out = "regulatory_tracker_pitch.pptx"
prs.save(out)
print(f"Saved: {out}")
